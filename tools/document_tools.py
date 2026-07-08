import os
from pathlib import Path
from tool_engine.tool_registry import register_tool, ToolPermission, ToolResult
from tools.file_tools_v2 import _validate_path



def _read_pdf(path: str) -> ToolResult:
    """读取 PDF 文件的文本和表格内容（前 20 页文本、前 5 页表格）。"""
    try:
        import pdfplumber
    except ImportError:
        return ToolResult.fail("PDF处理库未安装，请运行: pip install pdfplumber")

    try:
        target = os.path.expanduser(path)

        # 路径沙箱验证
        allowed, resolved, reason = _validate_path(target, mode="read")
        if not allowed:
            return ToolResult.fail(f"路径访问被拒绝: {reason}")

        if not os.path.exists(resolved):
            return ToolResult.fail(f"文件不存在: {path}")

        texts = []
        with pdfplumber.open(resolved) as pdf:
            total_pages = len(pdf.pages)
            for i, page in enumerate(pdf.pages[:20]):
                text = page.extract_text()
                if text:
                    texts.append(f"--- 第{i+1}页 ---\n{text}")

            tables = []
            for i, page in enumerate(pdf.pages[:5]):
                page_tables = page.extract_tables()
                for ti, table in enumerate(page_tables):
                    if table:
                        rows = [" | ".join(str(c or "") for c in row) for row in table[:10]]
                        tables.append(f"第{i+1}页 表格{ti+1}:\n" + "\n".join(rows))

        content = f"PDF: {resolved} ({total_pages}页)\n\n"
        content += "\n\n".join(texts[:3000])
        if tables:
            content += "\n\n--- 表格 ---\n" + "\n\n".join(tables[:1000])
        return ToolResult.ok(content[:5000])
    except Exception as e:
        return ToolResult.fail(f"PDF读取错误: {e!s}")


def _read_docx(path: str) -> ToolResult:
    """读取 DOCX 文件的段落文本和表格内容。"""
    try:
        from docx import Document
    except ImportError:
        return ToolResult.fail("DOCX处理库未安装，请运行: pip install python-docx")

    try:
        target = os.path.expanduser(path)

        # 路径沙箱验证
        allowed, resolved, reason = _validate_path(target, mode="read")
        if not allowed:
            return ToolResult.fail(f"路径访问被拒绝: {reason}")

        if not os.path.exists(resolved):
            return ToolResult.fail(f"文件不存在: {path}")

        doc = Document(resolved)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text for cell in row.cells]
                tables_text.append(" | ".join(cells))

        content = f"DOCX: {resolved}\n\n"
        content += "\n".join(paragraphs[:200])
        if tables_text:
            content += "\n\n--- 表格 ---\n" + "\n".join(tables_text[:50])
        return ToolResult.ok(content[:5000])
    except Exception as e:
        return ToolResult.fail(f"DOCX读取错误: {e!s}")


def _read_pptx(path: str) -> ToolResult:
    """读取 PPTX 文件每张幻灯片的文本内容（前 30 张）。"""
    try:
        from pptx import Presentation
    except ImportError:
        return ToolResult.fail("PPTX处理库未安装，请运行: pip install python-pptx")

    try:
        target = os.path.expanduser(path)

        # 路径沙箱验证
        allowed, resolved, reason = _validate_path(target, mode="read")
        if not allowed:
            return ToolResult.fail(f"路径访问被拒绝: {reason}")

        if not os.path.exists(resolved):
            return ToolResult.fail(f"文件不存在: {path}")

        prs = Presentation(resolved)
        slides_text = []
        for i, slide in enumerate(prs.slides[:30]):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
            if texts:
                slides_text.append(f"--- 幻灯片{i+1} ---\n" + "\n".join(texts))

        content = f"PPTX: {resolved} ({len(prs.slides)}页)\n\n"
        content += "\n\n".join(slides_text)
        return ToolResult.ok(content[:5000])
    except Exception as e:
        return ToolResult.fail(f"PPTX读取错误: {e!s}")


def _read_xlsx(path: str) -> ToolResult:
    """读取 XLSX 文件各工作表的行数据（前 5 个表，每表前 30 行）。"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ToolResult.fail("XLSX处理库未安装，请运行: pip install openpyxl")

    try:
        target = os.path.expanduser(path)

        # 路径沙箱验证
        allowed, resolved, reason = _validate_path(target, mode="read")
        if not allowed:
            return ToolResult.fail(f"路径访问被拒绝: {reason}")

        if not os.path.exists(resolved):
            return ToolResult.fail(f"文件不存在: {path}")

        wb = load_workbook(resolved, read_only=True, data_only=True)
        sheet_names = wb.sheetnames[:5]
        sheets_text = []
        for sheet_name in sheet_names:
            ws = wb[sheet_name]
            rows = []
            for i, row in enumerate(ws.iter_rows(max_row=30, values_only=True)):
                cells = [str(c) if c is not None else "" for c in row]
                rows.append(" | ".join(cells))
            if rows:
                sheets_text.append(f"--- {sheet_name} ---\n" + "\n".join(rows))

        sheet_count = len(wb.sheetnames)
        wb.close()
        content = f"XLSX: {resolved} ({sheet_count}个工作表)\n\n"
        content += "\n\n".join(sheets_text)
        return ToolResult.ok(content[:5000])
    except Exception as e:
        return ToolResult.fail(f"XLSX读取错误: {e!s}")


@register_tool(
    name="document_reader",
    description="读取文档内容。支持 PDF、DOCX、PPTX、XLSX 格式。输入文件路径。",
    schema={
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "文档文件路径"}
        },
        "required": ["path"],
    },
    permission=ToolPermission.READ_ONLY,
    category="document",
)
def document_reader(path: str) -> ToolResult:
    """根据文件扩展名分发到对应的文档读取器并返回内容。"""
    ext = Path(os.path.expanduser(path)).suffix.lower()
    readers = {
        '.pdf': _read_pdf,
        '.docx': _read_docx,
        '.doc': _read_docx,
        '.pptx': _read_pptx,
        '.ppt': _read_pptx,
        '.xlsx': _read_xlsx,
        '.xls': _read_xlsx,
    }
    reader = readers.get(ext)
    if not reader:
        supported = ", ".join(readers.keys())
        return ToolResult.fail(f"不支持的格式: {ext}。支持: {supported}")
    return reader(path)
